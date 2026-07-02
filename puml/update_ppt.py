#!/usr/bin/env python3
"""
update_ppt.py
Patches DRE architecture PPT to reflect the current architecture:
  - Genie Space → DQ Chatbot App (Databricks Apps + LangGraph)
  - Mosaic AI   → Model Serving Endpoint (self-hosted DBRX/BYOM)
  - Adds a new "LLM & Agent Stack on Databricks SaaS" clarification slide

Usage:
    python update_ppt.py
"""
from __future__ import annotations
import copy
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ──────────────────────────────────────────────────────────────────────
SRC  = Path(r"c:\Users\chattinr\OneDrive - adidas\Recordings\TDAPrj\drri_personal\tdarch-tda-data-monitoring\arch_v2\note-su\diagrams\DRE_Full_Architecture_All_Audiences_20260701_v2.pptx")
DEST = SRC.parent / "DRE_Full_Architecture_All_Audiences_20260702_v3.pptx"

# ── Text substitution map ──────────────────────────────────────────────────────
REPLACEMENTS: list[tuple[str, str]] = [
    # Genie Space variants → DQ Chatbot App
    ("Genie Space / Agent Streamlit UI",  "DQ Chatbot App (Databricks Apps + LangGraph)"),
    ("Agent Streamlit UI",                "DQ Chatbot App (LangGraph + Streamlit)"),
    ("Genie Space [Chat UI]",             "DQ Chatbot App [LangGraph + Streamlit UI]"),
    ("Genie Space / Agent Streamlit",     "DQ Chatbot App (LangGraph)"),
    ("Genie Space",                       "DQ Chatbot App (Databricks Apps)"),
    ("Genie checks",                      "DQ Chatbot App checks"),
    ("Genie asks",                        "DQ Chatbot App asks"),
    ("Genie proposes",                    "DQ Chatbot App proposes"),
    ("Genie Suggestion Cache Strategy",   "DQ Chatbot App — Rule Suggestion Cache"),
    ("Genie Suggestion Cache",            "Rule Suggestion Cache (Delta table)"),
    ("live in Genie Space",               "via DQ Chatbot App"),
    ("Genie + LLM",                       "DQ Chatbot App + Model Serving Endpoint"),
    ("Genie Space query",                 "DQ Chatbot App query"),
    ("Genie Space and",                   "DQ Chatbot App and"),

    # Mosaic AI → Model Serving Endpoint
    ("Mosaic AI / Claude",                "Model Serving Endpoint (DBRX / BYOM, self-hosted)"),
    ("Mosaic AI | Databricks Serverless", "Model Serving Endpoint (Serverless inference)"),
    ("Databricks Workspace | Mosaic AI",  "Databricks Workspace | Model Serving Endpoint"),
    ("Diagnosis + Deflection | Mosaic AI","Diagnosis + Deflection | Model Serving Endpoint"),
    ("Mosaic AI / LLM Service Endpoint",  "Model Serving Endpoint (DBRX / BYOM)"),
    ("Mosaic AI",                         "Model Serving Endpoint"),

    # Generic Genie leftover
    ("via Genie Space",                   "via DQ Chatbot App"),
    ("in Genie Space",                    "in DQ Chatbot App"),
    ("through Genie Space",               "through DQ Chatbot App"),
]


def patch_run(run_text: str) -> tuple[str, bool]:
    """Apply all substitutions to a single text run. Returns (new_text, changed)."""
    changed = False
    for old, new in REPLACEMENTS:
        if old in run_text:
            run_text = run_text.replace(old, new)
            changed = True
    return run_text, changed


def patch_presentation(prs: Presentation) -> int:
    """Walk all text frames and apply substitutions. Returns number of changes."""
    total = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    new_text, changed = patch_run(run.text)
                    if changed:
                        run.text = new_text
                        total += 1
    return total


def add_saas_clarification_slide(prs: Presentation) -> None:
    """Insert a new slide after the current 'Platform & Component Map' slide
    that clearly maps each DRE component to its Databricks SaaS service."""

    # Find the "Platform & Component Map" slide (slide 28 in original)
    target_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if "Platform & Component Map" in para.text:
                        target_idx = i
                        break

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank
    new_slide = prs.slides.add_slide(blank_layout)

    W = prs.slide_width
    H = prs.slide_height

    def add_text_box(slide, left, top, width, height, text,
                     font_size=16, bold=False, color=(0x1A, 0x23, 0x7E),
                     fill_rgb=None, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
        return txBox

    def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
        from pptx.util import Pt
        shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
        if line_rgb:
            shape.line.color.rgb = RGBColor(*line_rgb)
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    # Title
    add_text_box(new_slide, Inches(0.3), Inches(0.1), W - Inches(0.6), Inches(0.6),
                 "LangGraph · RAG KB · LLM — Hosted on Which Databricks SaaS?",
                 font_size=20, bold=True, color=(0x1A, 0x23, 0x7E),
                 align=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(new_slide, Inches(0.3), Inches(0.65), W - Inches(0.6), Inches(0.35),
                 "Every DRE component runs on a managed Databricks SaaS service — no extra infrastructure.",
                 font_size=12, bold=False, color=(0x55, 0x55, 0x55),
                 align=PP_ALIGN.CENTER)

    # Table of components
    rows = [
        # (Component, Databricks SaaS, Purpose, colour_fill, colour_text)
        ("DQ Chatbot App\n(Streamlit UI)",
         "Databricks Apps (Serverless)",
         "Hosts the LangGraph runtime and Streamlit chat UI.\nSSO / OAuth · Auto-scale · Port 443.",
         (0xE8, 0xEA, 0xF6), (0x1A, 0x23, 0x7E)),

        ("LangGraph Orchestrator",
         "Runs inside Databricks Apps",
         "State graph manages tool calls, intent routing,\nmemory load/save and CRAG retrieval per turn.",
         (0xE3, 0xF2, 0xFD), (0x0D, 0x47, 0xA1)),

        ("LLM Inference\n(SQL gen · root-cause · NL)",
         "Model Serving Endpoint\n(DBRX / BYOM — self-hosted)",
         "Serverless endpoint · autoscale → 0 · <5 s p95\nTemperature 0.1 for SQL, 0.7 for NL suggestions.",
         (0xFF, 0xF8, 0xE1), (0xE6, 0x51, 0x00)),

        ("AI Gateway",
         "Databricks AI Gateway",
         "Rate limits · token budgets per team · PII filter\nCost caps · request routing to Model Serving.",
         (0xFF, 0xF3, 0xE0), (0xBF, 0x36, 0x0C)),

        ("RAG Knowledge Base\n(doc chunks)",
         "UC Volume + Databricks Vector Search",
         "Raw docs in UC Volume (data_contracts/, sql_patterns/).\nBGE-M3 embeddings · CRAG confidence gate (≥0.72 use).",
         (0xE8, 0xF5, 0xE9), (0x1B, 0x5E, 0x20)),

        ("Delta Lake / Unity Catalog\n(all state tables)",
         "Delta Lake governed by Unity Catalog",
         "All 6 state tables registered in UC.\nACID · time-travel audit · row/column security.",
         (0xE0, 0xF7, 0xFA), (0x00, 0x60, 0x64)),

        ("D1 / D2 / D3 Agent Jobs",
         "Databricks Workflows (Serverless jobs)",
         "D1 scheduled + event trigger · D2 on-demand · D3 triggered.\nMulti-task job DAGs · Delta trigger integration.",
         (0xF3, 0xE5, 0xF5), (0x6A, 0x1B, 0x9A)),

        ("Lakehouse Monitoring",
         "Databricks Lakehouse Monitoring",
         "EWMA baselines · token cost dashboards\nAnomaly spike alert (<5 min) · weekly domain report.",
         (0xEF, 0xEB, 0xE9), (0x4E, 0x34, 0x2E)),
    ]

    row_h = Inches(0.55)
    col_widths = [Inches(2.1), Inches(2.4), Inches(4.5)]
    left_start = Inches(0.3)
    top_start = Inches(1.1)

    # Header row
    header_fills = [(0x37, 0x47, 0x4F), (0x37, 0x47, 0x4F), (0x37, 0x47, 0x4F)]
    header_texts = ["DRE Component", "Databricks SaaS Service", "Purpose & Key Config"]
    for ci, (htxt, hfill) in enumerate(zip(header_texts, header_fills)):
        cx = left_start + sum(col_widths[:ci])
        add_rect(new_slide, cx, top_start, col_widths[ci], Inches(0.38),
                 hfill, (0x27, 0x37, 0x4F))
        add_text_box(new_slide, cx + Inches(0.05), top_start + Inches(0.04),
                     col_widths[ci] - Inches(0.1), Inches(0.38),
                     htxt, font_size=11, bold=True,
                     color=(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    # Data rows
    for ri, (comp, svc, purpose, fill_rgb, text_rgb) in enumerate(rows):
        row_top = top_start + Inches(0.38) + ri * row_h
        row_data = [comp, svc, purpose]
        for ci, cell_txt in enumerate(row_data):
            cx = left_start + sum(col_widths[:ci])
            add_rect(new_slide, cx, row_top, col_widths[ci], row_h,
                     fill_rgb if ci < 2 else (0xFF, 0xFF, 0xFF),
                     (0xBD, 0xBD, 0xBD))
            add_text_box(new_slide, cx + Inches(0.05), row_top + Inches(0.03),
                         col_widths[ci] - Inches(0.1), row_h - Inches(0.06),
                         cell_txt,
                         font_size=9 if ci == 2 else 10,
                         bold=(ci == 0),
                         color=text_rgb if ci < 2 else (0x33, 0x33, 0x33))

    # Footer note
    footer_top = top_start + Inches(0.38) + len(rows) * row_h + Inches(0.1)
    add_text_box(new_slide, Inches(0.3), footer_top, W - Inches(0.6), Inches(0.4),
                 "★  All services run inside the same Databricks workspace · All data stays in Unity Catalog · No external infra required",
                 font_size=10, bold=True, color=(0x1B, 0x5E, 0x20),
                 align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print(f"Loading: {SRC.name}")
    prs = Presentation(str(SRC))

    changes = patch_presentation(prs)
    print(f"Text replacements applied: {changes}")

    add_saas_clarification_slide(prs)
    print("Added new slide: 'LangGraph · RAG KB · LLM — Hosted on Which Databricks SaaS?'")

    prs.save(str(DEST))
    print(f"\n✅  Saved: {DEST.name}")
    print(f"   ({len(prs.slides)} slides total)")
